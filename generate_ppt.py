#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI+PLC 产品规划 PPT 生成器
基于 plan.md 内容生成专业的商业计划汇报演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """创建并返回完整的演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 设置主题色
    PRIMARY_COLOR = RGBColor(0, 102, 204)  # 蓝色
    ACCENT_COLOR = RGBColor(255, 102, 0)   # 橙色
    DARK_GRAY = RGBColor(64, 64, 64)
    LIGHT_GRAY = RGBColor(128, 128, 128)
    
    def add_title_slide(title, subtitle):
        """添加标题幻灯片"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        
        # 背景色块
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 248, 255)
        shape.line.fill.background()
        
        # 标题
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        left = Inches(1)
        top = Inches(4.2)
        width = Inches(8)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, bullets, slide_number=None):
        """添加内容幻灯片（带项目符号）"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 标题区域
        left = Inches(0.5)
        top = Inches(0.4)
        width = Inches(9)
        height = Inches(0.8)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # 添加标题下划线
        left = Inches(0.5)
        top = Inches(1.1)
        width = Inches(9)
        height = Inches(0.02)
        line_shape = slide.shapes.add_shape(1, left, top, width, height)
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = ACCENT_COLOR
        line_shape.line.fill.background()
        
        # 内容区域
        left = Inches(0.7)
        top = Inches(1.5)
        width = Inches(8.6)
        height = Inches(5.5)
        
        text_frame = slide.shapes.add_textbox(left, top, width, height).text_frame
        text_frame.word_wrap = True
        
        for i, bullet_text in enumerate(bullets):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            
            # 判断层级（通过缩进）
            level = 0
            clean_text = bullet_text
            if bullet_text.startswith('  - '):
                level = 1
                clean_text = bullet_text[4:]
            elif bullet_text.startswith('    '):
                level = 2
                clean_text = bullet_text[4:]
            elif bullet_text.startswith('- '):
                level = 0
                clean_text = bullet_text[2:]
            
            p.text = clean_text
            p.level = level
            p.font.size = Pt(16) if level == 0 else Pt(14)
            p.font.color.rgb = DARK_GRAY if level == 0 else LIGHT_GRAY
            p.space_after = Pt(8)
            
            # 加粗关键词
            if '：' in clean_text or '：' in clean_text:
                p.font.bold = True
        
        # 幻灯片编号
        if slide_number:
            left = Inches(9.2)
            top = Inches(7)
            width = Inches(0.5)
            height = Inches(0.3)
            num_box = slide.shapes.add_textbox(left, top, width, height)
            tf = num_box.text_frame
            tf.text = str(slide_number)
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = LIGHT_GRAY
            p.alignment = PP_ALIGN.RIGHT
        
        return slide
    
    # 封面
    add_title_slide(
        "AI+PLC：下一代工业编程操作系统",
        "技术可行性与商业价值分析 | 2025年度战略项目"
    )
    
    # Slide 1: 战略定位
    add_content_slide(
        "Slide 1 | AI+PLC：下一条可验证的增长曲线",
        [
            "汇报目的：争取 12 个月、1200 万人民币的专项投入，押注「AI工业编程」作为公司下一条增长曲线",
            "",
            "行业信号：2024 以来，Siemens、ABB 等都发表 PLC+LLM 研究（LLM4PLC、Agents4PLC、AutoPLC），表明技术从「概念验证」进入「验证规模化」阶段",
            "",
            "我们优势：掌握真实产线数据+行业 know-how，可把一线工程师的重复编码沉淀成可复用的软件资产，形成跨品牌护城河",
            "",
            "技术成熟度：学术界已验证可控、可验证、可商用的技术路线，编译通过率达 90%+，推理成本仅 $0.13/任务"
        ],
        1
    )
    
    # Slide 2: 市场需求
    add_content_slide(
        "Slide 2 | 客户拉力：OEM 与集成商在主动催进度",
        [
            "市场信号：CEChina（2025.04）显示 OEM/系统集成商/分销商把生成式 AI 视作核心工具，Beckhoff TwinCAT Chat 2023 展会被抢着报名验证",
            "",
            "一线痛点：模板化标签、批量 POU、文档、调试占 >60% 工时；人才断层导致客户愿意用 AI 加速新人上岗",
            "",
            "时间窗口：厂商自研 Copilot（Siemens、Rockwell）锁定自家生态，12–18 个月内跨品牌解决方案仍空白，我们必须抢占心智",
            "",
            "竞争态势：Siemens Industrial Copilot 已发布但仅支持 SCL/TIA Portal，缺乏 Ladder/SFC 与跨品牌能力"
        ],
        2
    )
    
    # Slide 3: 技术验证
    add_content_slide(
        "Slide 3 | 技术验证：三条路线共同证明「可控、可验证、可商用」",
        [
            "LLM4PLC（ICSE-SEIP 2024）：引入用户反馈+语法检查+编译器+SMV 模型检验，配合 LoRA 微调，在 Fischer Technik 产线测试中把生成成功率从 47% 拉到 72%，专家评分 2.25→7.75/10",
            "",
            "Agents4PLC（arXiv:2410.14209）：首个多智能体闭环，从需求→规划→编码→编译→形式验证→调试，强调「代码级」验证，可插拔 GPT、CodeLlama 等模型，还发布了从自然语言到 ST 代码的首个基准",
            "",
            "AutoPLC（arXiv:2412.02410）：构建「厂商 API 库 + 案例库 + IDE 动态验证」四阶段流水线，在 Siemens TIA Portal/CODESYS 914 个任务上实现 90%+ 编译通过，单任务推理成本 0.13 美元且获得资深工程师认可",
            "",
            "结论：技术路径已被国际顶会和头部厂商验证，风险可控"
        ],
        3
    )
    
    # Slide 4: 差异化定位
    add_content_slide(
        "Slide 4 | 机会缺口：差异化聚焦 Ladder/SFC 与跨厂商落地",
        [
            "现状分析：现有方案几乎清一色聚焦 Structured Text；AutoPLC 也强调不同厂商 SCL/ST 差异巨大，客户需要「懂设备方言」的 AI",
            "",
            "Siemens Industrial Copilot 局限性：",
            "  - 功能：仅支持 SCL 代码生成与 WinCC Unified 可视化，锁定 Simatic 生态",
            "  - 依赖：需要 Azure OpenAI 订阅、WinCC Unified 许可，本地安装约 329 MB",
            "  - 可用性：目前仅对欧美少量客户开放订阅，中国客户仍需等待",
            "",
            "我们的差异化优势：",
            "  - 跨品牌支持：Siemens + CODESYS + Rockwell + Beckhoff 全栈覆盖",
            "  - Ladder/SFC 增强：提供文本中间表示、ASCII 渲染与回写，比只生成 SCL 更贴合现场遗留项目",
            "  - 私有化部署：可插拔推理源，支持国内算力与完全离线环境"
        ],
        4
    )
    
    # Slide 5: 产品蓝图
    add_content_slide(
        "Slide 5 | 产品蓝图：AI 工业编程操作系统",
        [
            "知识底座：行业标准（IEC 61131-3/61499）+ 厂商文档（API/错误码/最佳实践）+ 真实案例库（已验证的产线代码）",
            "",
            "双层模型：通用大模型（GPT-4/国产基座）+ 领域微调模型（LoRA on 行业数据集）",
            "",
            "Multi-Agent 工作流：需求分析 → 架构规划 → 代码生成 → 语法检查 → 编译验证 → 形式验证（SMV）→ 仿真测试 → 文档生成",
            "",
            "RAG + 守护进程：实时检索厂商手册、案例库、调试日志，持续监控生成质量并触发人工审核",
            "",
            "与竞品差异化设计：",
            "  - 借鉴 Siemens Copilot 的「IDE 插件 + 云模型」模式，但支持可插拔推理源与私有化",
            "  - 增强 Ladder/SFC：文本中间表示、ASCII 渲染与回写",
            "  - 多厂商扩展：真正的「工业多语种 Copilot」"
        ],
        5
    )
    
    # Slide 6: 商业模式
    add_content_slide(
        "Slide 6 | 商业模式与单位经济假设",
        [
            "客户分层：",
            "  - ① 大型 OEM/系统集成商（席位+算力包）",
            "  - ② PLC 厂商（白标/SDK，共建生态）",
            "  - ③ 终端工厂（项目制+成功费）",
            "",
            "收费建议：基础席位 30–50 万元/年 + 生成量计费；AutoPLC 报告 $0.13/任务的推理成本，给我们定价与毛利守护空间",
            "",
            "ROI 案例：CEChina 案例显示重复编码时间可砍半，典型 4 人月项目可节省 20–30 万元/线；若叠加私有化与安全白盒审计，可再收一次性部署费",
            "",
            "生态协同：与分销商（DigiKey 等）共建应用库，与高校/培训机构共建联合实验室，持续获得长尾设备与数据"
        ],
        6
    )
    
    # Slide 7: 执行计划
    add_content_slide(
        "Slide 7 | 12 个月执行里程碑（对齐论文验证路径）",
        [
            "Q1：签 3 家灯塔客户，清洗 ≥300 个 ST 项目，搭建基础 RAG+验证流水线，完成封闭环境 PoC（借鉴 LLM4PLC 的语法/SMV 流程）",
            "",
            "Q2：上线「需求→ST 代码→IDE 验证」MVP（Siemens+Codesys 双栈），同步完成 Ladder/SFC 文本中间表示与 ASCII 渲染 Demo",
            "",
            "Q3：引入 Agents4PLC 式多 Agent 协同（诊断、文档、测试自动生成），交付首批付费项目，启动私有化/合规模块",
            "",
            "Q4：发布 2.0（支持 HMI、SFC、知识问答），建立行业基准测试与合作伙伴认证，准备下一轮融资/规模化推广",
            "",
            "资源需求：核心团队 12 人（控制算法 4、LLM 4、后端 2、产品/交付 2），年度研发+算力预算约 1200 万人民币"
        ],
        7
    )
    
    # Slide 8: 风险应对
    add_content_slide(
        "Slide 8 | 主要风险与应对（结合论文经验）",
        [
            "幻觉与质量：全链路必须包含编译/仿真/SMV/单测，关键逻辑强制人工签核；沿用 LLM4PLC、Agents4PLC 的多级验证策略",
            "",
            "数据安全：提供本地部署、脱敏管道、权限审计；AutoPLC 的厂商协作案例说明 OEM 接受「私有模型+本地 IDE」模式",
            "",
            "合规与可解释性：跟进 EU AI Act、国内等保，保留模型决策日志与调试记录，支撑审计",
            "",
            "竞争：厂商 Copilot 绑定单品牌；我们主打「跨品牌 + Ladder/SFC + 私有化」，并与器件分销商、高校共建生态提高进入壁垒",
            "",
            "人才：联合高校/行业协会建立实验室，提前培育「控制+LLM」复合型工程师"
        ],
        8
    )
    
    # Slide 9: 决策请求
    add_content_slide(
        "Slide 9 | 管理层需拍板事项",
        [
            "决策请求：",
            "  - ① 是否批准 1200 万人民币预算与 12 人编制",
            "  - ② 是否确认 3 家灯塔客户名单与数据采集协议",
            "  - ③ 是否授权启动 Q1 PoC 与 Ladder/SFC 中间表示研发",
            "",
            "下一步：若获批，即刻按照 Q1 计划启动数据治理、PoC 与技术验证",
            "",
            "信息来源：",
            "  - LLM4PLC（ICSE-SEIP 2024）",
            "  - Agents4PLC（arXiv:2410.14209）",
            "  - AutoPLC（arXiv:2412.02410）",
            "  - IEC 图形语言探索（arXiv:2410.15200）",
            "  - CEChina 2025、Beckhoff 2023"
        ],
        9
    )
    
    # 附录：技术路线
    add_content_slide(
        "附录 | Ladder/SFC 技术路线（基于论文实验）",
        [
            "实验发现：GPT-4 在提供少量示例后可生成可读 SFC，但 Ladder 受 ASCII 布局/连线约束影响大，必须拆解成中间表示再渲染",
            "",
            "可执行路径：",
            "  - ① 建立 JSON/文本中间表示 → 生成后再渲染成 Ladder 图",
            "  - ② 对常见梯级（互锁、延时、安全链）做库化，AI 只拼装已验证模块",
            "  - ③ 用 RAG+图语法限制触点/线圈幻觉",
            "  - ④ 结合 Siemens PLCSIM/CODESYS 仿真闭环验证，人工终审",
            "",
            "技术门槛：需要深度理解 IEC 61131-3 标准与各厂商的 Ladder 方言差异，但已有明确的工程化路径"
        ],
        10
    )
    
    return prs


def main():
    """主函数"""
    print("开始生成 AI+PLC 产品规划 PPT...")
    
    prs = create_presentation()
    
    output_file = "/home/engine/project/AI+PLC产品规划汇报.pptx"
    prs.save(output_file)
    
    print(f"✅ PPT 生成成功: {output_file}")
    print(f"📊 包含 {len(prs.slides)} 张幻灯片")
    print("\n幻灯片列表:")
    print("  - 封面：AI+PLC：下一代工业编程操作系统")
    print("  - Slide 1：战略定位与增长曲线")
    print("  - Slide 2：市场需求与客户拉力")
    print("  - Slide 3：技术验证（三篇论文）")
    print("  - Slide 4：差异化定位与竞争优势")
    print("  - Slide 5：产品架构蓝图")
    print("  - Slide 6：商业模式与单位经济")
    print("  - Slide 7：12 个月执行里程碑")
    print("  - Slide 8：风险应对策略")
    print("  - Slide 9：管理层决策请求")
    print("  - 附录：Ladder/SFC 技术路线")


if __name__ == "__main__":
    main()
